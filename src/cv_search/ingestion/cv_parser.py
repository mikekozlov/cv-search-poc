# src/cvsearch/cv_parser.py
from __future__ import annotations
from pathlib import Path
from typing import List
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph


class CVParser:
    """
    Handles the extraction of raw text from .pptx CV files.
    """

    def extract_text(self, file_path: Path) -> str:
        """
        Opens a .pptx file and extracts all text content, preserving
        paragraphs with newline separators.

        Args:
            file_path: The Path object pointing to the .pptx file.

        Returns:
            A single string containing all extracted text.
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found at: {file_path}")

        try:
            prs = Presentation(file_path)
            text_runs: List[str] = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not self._has_text_frame(shape):
                        continue

                    for paragraph in shape.text_frame.paragraphs:
                        text_runs.append(self._get_paragraph_text(paragraph))

                text_runs.append("\n")

            return "\n".join(text_runs).strip()

        except Exception as e:
            print(f"Error processing {file_path.name}: {e}")
            raise

    def _has_text_frame(self, shape: Shape) -> bool:
        return hasattr(shape, "text_frame") and shape.text_frame is not None

    def _get_paragraph_text(self, paragraph: _Paragraph) -> str:
        return "".join(run.text for run in paragraph.runs)
